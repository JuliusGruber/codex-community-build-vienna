"""Generate the Codex Community Build - Vienna deck.

Styling mirrors https://codex-events.com/ : white canvas, near-black ink
(#030213), a dark circular logo mark, dark uppercase letter-spaced tag pills,
black pill CTAs and a thin muted footer rule. Run:

    python slides/build_slides.py

Produces slides/codex-community-build-vienna.pptx
"""

from pathlib import Path

from lxml import etree
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.opc.constants import RELATIONSHIP_TYPE as RT
from pptx.util import Emu, Inches, Pt

# --- palette (sampled from codex-events.com) -------------------------------
INK = RGBColor(0x03, 0x02, 0x13)      # primary text / brand dark
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
MUTED = RGBColor(0x6E, 0x6E, 0x80)    # secondary text
FAINT = RGBColor(0x9A, 0x9A, 0xA8)    # footer / fine print
LINE = RGBColor(0xE6, 0xE6, 0xEB)     # hairline rules

FONT = "Segoe UI"
FONT_SEMI = "Segoe UI Semibold"

EMU_IN = 914400
SLIDE_W = 13.333
SLIDE_H = 7.5
MARGIN = 0.9
CONTENT_W = SLIDE_W - 2 * MARGIN


def _set_run(run, text, *, size, color=INK, font=FONT, bold=False, spc=None):
    run.text = text
    f = run.font
    f.size = Pt(size)
    f.name = font
    f.bold = bold
    f.color.rgb = color
    if spc is not None:  # letter spacing in 1/100 pt
        run._r.get_or_add_rPr().set("spc", str(int(spc)))


def _no_autofit(tf):
    tf.word_wrap = True
    # MSO_AUTO_SIZE.NONE == 0
    try:
        from pptx.enum.text import MSO_AUTO_SIZE

        tf.auto_size = MSO_AUTO_SIZE.NONE
    except Exception:
        pass


def set_link_theme_color(prs, hex6="030213"):
    """PowerPoint renders hyperlink text in the THEME hlink color, ignoring the
    run's explicit color. Override hlink / folHlink in every theme to ink."""
    a = "http://schemas.openxmlformats.org/drawingml/2006/main"
    for master in prs.slide_masters:
        part = master.part.part_related_by(RT.THEME)
        root = etree.fromstring(part.blob)
        changed = False
        for tag in ("hlink", "folHlink"):
            el = root.find(f".//{{{a}}}{tag}")
            if el is None:
                continue
            for child in list(el):
                el.remove(child)
            etree.SubElement(el, f"{{{a}}}srgbClr", {"val": hex6})
            changed = True
        if changed:
            part._blob = etree.tostring(root, xml_declaration=True, encoding="UTF-8", standalone=True)


def add_background(slide):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = WHITE


def add_logo(slide):
    """Dark circle mark + small wordmark, top-left (matches the site header)."""
    d = 0.42
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.62), Inches(0.46), Inches(d), Inches(d))
    circle.fill.solid()
    circle.fill.fore_color.rgb = INK
    circle.line.fill.background()
    # inner ring for a touch of mark detail
    inner = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, Inches(0.62 + d * 0.30), Inches(0.46 + d * 0.30), Inches(d * 0.40), Inches(d * 0.40)
    )
    inner.fill.background()
    inner.line.color.rgb = WHITE
    inner.line.width = Pt(1.25)

    box = slide.shapes.add_textbox(Inches(1.18), Inches(0.44), Inches(4.5), Inches(0.5))
    tf = box.text_frame
    _no_autofit(tf)
    tf.margin_left = tf.margin_top = tf.margin_bottom = 0
    p1 = tf.paragraphs[0]
    p1.line_spacing = 1.0
    _set_run(p1.add_run(), "CODEX COMMUNITY", size=9.5, color=INK, font=FONT_SEMI, spc=140)
    p2 = tf.add_paragraph()
    p2.line_spacing = 1.0
    _set_run(p2.add_run(), "Build · Vienna", size=9, color=MUTED, font=FONT)


def add_footer(slide, page_no):
    rule = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(MARGIN), Inches(6.92), Inches(CONTENT_W), Pt(1)
    )
    rule.fill.solid()
    rule.fill.fore_color.rgb = LINE
    rule.line.fill.background()

    box = slide.shapes.add_textbox(Inches(MARGIN), Inches(7.0), Inches(CONTENT_W), Inches(0.35))
    tf = box.text_frame
    _no_autofit(tf)
    tf.margin_left = tf.margin_top = 0
    p = tf.paragraphs[0]
    _set_run(p.add_run(), "codex-events.com   ·   Codex Community Build — Vienna", size=8.5, color=FAINT, font=FONT)
    r = p.add_run()
    _set_run(r, f"        {page_no}", size=8.5, color=FAINT, font=FONT)
    # page number right-aligned in its own box
    pn = slide.shapes.add_textbox(Inches(SLIDE_W - MARGIN - 1.0), Inches(7.0), Inches(1.0), Inches(0.35))
    ptf = pn.text_frame
    _no_autofit(ptf)
    pp = ptf.paragraphs[0]
    pp.alignment = PP_ALIGN.RIGHT
    _set_run(pp.add_run(), f"{page_no:02d}", size=8.5, color=FAINT, font=FONT_SEMI, spc=60)
    # blank the placeholder page text in footer box
    p.clear()
    _set_run(p.add_run(), "codex-events.com   ·   Codex Community Build — Vienna", size=8.5, color=FAINT, font=FONT)


def add_tag(slide, text, left, top):
    w = max(0.7, 0.115 * len(text) + 0.42)
    pill = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(w), Inches(0.32))
    pill.adjustments[0] = 0.5
    pill.fill.solid()
    pill.fill.fore_color.rgb = INK
    pill.line.fill.background()
    tf = pill.text_frame
    tf.margin_left = tf.margin_right = Inches(0.08)
    tf.margin_top = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    _set_run(p.add_run(), text.upper(), size=9, color=WHITE, font=FONT_SEMI, spc=160)
    return w


def add_heading(slide, text, lead=None, top=1.5):
    box = slide.shapes.add_textbox(Inches(MARGIN), Inches(top), Inches(CONTENT_W), Inches(0.8))
    tf = box.text_frame
    _no_autofit(tf)
    tf.margin_left = tf.margin_top = 0
    p = tf.paragraphs[0]
    _set_run(p.add_run(), text, size=30, color=INK, font=FONT_SEMI)
    # accent bar
    bar = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(MARGIN), Inches(top + 0.72), Inches(0.62), Inches(0.06))
    bar.adjustments[0] = 0.5
    bar.fill.solid()
    bar.fill.fore_color.rgb = INK
    bar.line.fill.background()
    if lead:
        lb = slide.shapes.add_textbox(Inches(MARGIN), Inches(top + 0.92), Inches(CONTENT_W), Inches(0.5))
        ltf = lb.text_frame
        _no_autofit(ltf)
        ltf.margin_left = ltf.margin_top = 0
        lp = ltf.paragraphs[0]
        _set_run(lp.add_run(), lead, size=14, color=MUTED, font=FONT)


def add_bullets(slide, items, *, left=MARGIN, top=2.7, width=CONTENT_W, size=16, gap=9):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(6.7 - top))
    tf = box.text_frame
    _no_autofit(tf)
    tf.margin_left = tf.margin_top = 0
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.line_spacing = 1.12
        p.space_after = Pt(gap)
        _set_run(p.add_run(), "—  ", size=size, color=INK, font=FONT_SEMI)
        if isinstance(item, tuple):
            head, rest = item
            _set_run(p.add_run(), head, size=size, color=INK, font=FONT_SEMI)
            if rest:
                _set_run(p.add_run(), rest, size=size, color=MUTED, font=FONT)
        else:
            _set_run(p.add_run(), item, size=size, color=INK, font=FONT)
    return box


def base_slide(prs, page_no, *, logo=True, footer=True):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    if logo:
        add_logo(slide)
    if footer:
        add_footer(slide, page_no)
    return slide


def build():
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)

    # --- 1. Title ----------------------------------------------------------
    s = base_slide(prs, 1)
    add_tag(s, "Build", MARGIN, 2.55)
    box = s.shapes.add_textbox(Inches(MARGIN), Inches(3.0), Inches(CONTENT_W), Inches(1.6))
    tf = box.text_frame
    _no_autofit(tf)
    tf.margin_left = tf.margin_top = 0
    p = tf.paragraphs[0]
    p.line_spacing = 1.02
    _set_run(p.add_run(), "Codex Community Build", size=46, color=INK, font=FONT_SEMI)
    p2 = tf.add_paragraph()
    p2.line_spacing = 1.02
    _set_run(p2.add_run(), "Vienna", size=46, color=INK, font=FONT_SEMI)
    sub = s.shapes.add_textbox(Inches(MARGIN), Inches(4.95), Inches(CONTENT_W), Inches(0.6))
    stf = sub.text_frame
    _no_autofit(stf)
    stf.margin_left = stf.margin_top = 0
    sp = stf.paragraphs[0]
    _set_run(sp.add_run(), "20 June 2026", size=15, color=INK, font=FONT_SEMI)
    _set_run(sp.add_run(), "   ·   Vienna, Austria", size=15, color=MUTED, font=FONT)

    # --- 2. Links (just the two links) ------------------------------------
    s = base_slide(prs, 2)
    kb = s.shapes.add_textbox(Inches(MARGIN), Inches(1.7), Inches(CONTENT_W), Inches(0.4))
    ktf = kb.text_frame
    _no_autofit(ktf)
    ktf.margin_left = ktf.margin_top = 0
    kp = ktf.paragraphs[0]
    _set_run(kp.add_run(), "PROJECT REPOSITORIES", size=11, color=MUTED, font=FONT_SEMI, spc=180)
    links = [
        "https://github.com/JuliusGruber/codex-community-build-vienna",
        "https://github.com/JuliusGruber/AI-Engineering-Coach",
    ]
    y = 2.55
    for url in links:
        # black pill marker
        dot = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(MARGIN), Inches(y + 0.05), Inches(0.34), Inches(0.34))
        dot.adjustments[0] = 0.5
        dot.fill.solid()
        dot.fill.fore_color.rgb = INK
        dot.line.fill.background()
        atf = dot.text_frame
        atf.margin_left = atf.margin_right = atf.margin_top = atf.margin_bottom = 0
        ap = atf.paragraphs[0]
        ap.alignment = PP_ALIGN.CENTER
        _set_run(ap.add_run(), "→", size=13, color=WHITE, font=FONT_SEMI)
        # link text
        lb = s.shapes.add_textbox(Inches(MARGIN + 0.55), Inches(y), Inches(CONTENT_W - 0.55), Inches(0.6))
        ltf = lb.text_frame
        _no_autofit(ltf)
        ltf.vertical_anchor = MSO_ANCHOR.MIDDLE
        ltf.margin_left = ltf.margin_top = ltf.margin_bottom = 0
        lp = ltf.paragraphs[0]
        run = lp.add_run()
        run.hyperlink.address = url  # set link first so the explicit color/underline below win
        _set_run(run, url.replace("https://", ""), size=24, color=INK, font=FONT_SEMI)
        run.font.underline = False
        y += 1.15

    # --- 3. AGENTS.md ------------------------------------------------------
    s = base_slide(prs, 3)
    add_heading(s, "AGENTS.md", lead="Durable, layered project guidance Codex loads from repo root toward the working directory.")
    add_bullets(
        s,
        [
            ("Layered files — ", "global ~/.codex/, repo-root AGENTS.md, nested AGENTS.md / AGENTS.override.md"),
            ("Closer overrides broader — ", "one file per directory; ~32 KiB combined limit"),
            ("A good one covers — ", "layout, how to run, build/test/lint, conventions & PR expectations, do-not rules, what “done” means"),
            ("Keep it practical — ", "short and accurate beats long and vague"),
            ("/init ", "scaffolds a starter; update it when Codex repeats a mistake"),
            ("Open format — ", "stewarded by the Agentic AI Foundation (Linux Foundation)"),
        ],
        top=2.75,
        gap=10,
    )

    # --- 4. Codex Best Practices (section headings of wiki.md) -------------
    s = base_slide(prs, 4)
    add_heading(s, "Codex Best Practices", lead="Treat Codex like a teammate you configure and improve over time.")
    headings = [
        "Speech dictation",
        "Model and reasoning effort (CLI)",
        "Configure the status line (CLI)",
        "Configure Codex",
        "Enable Live Web Search",
        "skills.md",
        "MCP",
        "Worktrees",
        "Plan mode",
        "Prompt essentials",
        "Review",
        "Interrogatory LLM",
    ]
    col_w = CONTENT_W / 2 - 0.2
    add_bullets(s, headings[:6], left=MARGIN, top=2.8, width=col_w, size=16, gap=11)
    add_bullets(s, headings[6:], left=MARGIN + CONTENT_W / 2 + 0.2, top=2.8, width=col_w, size=16, gap=11)

    # --- 5. Top MCP servers ------------------------------------------------
    s = base_slide(prs, 5)
    add_heading(s, "Top MCP servers for Codex", lead="Most-used MCP servers (ranked); 1–6 are high-confidence — named in OpenAI’s official Codex docs.")
    add_bullets(
        s,
        [
            ("Context7 — ", "up-to-date, version-pinned library docs"),
            ("Playwright MCP — ", "browser automation via accessibility-tree snapshots"),
            ("GitHub MCP — ", "PRs, issues, code search, CI logs"),
            ("Chrome DevTools MCP — ", "drives live Chrome via CDP (perf, DOM, network)"),
            ("Sentry MCP — ", "production errors & traces for debugging"),
            ("Figma MCP — ", "design specs/tokens → code (hosted)"),
            ("Supabase MCP — ", "Postgres/SQL, migrations, auth/storage"),
            ("Also — ", "Filesystem · Serena · OpenAI Docs MCP"),
        ],
        top=2.75,
        size=15.5,
        gap=7,
    )

    # --- 6. Common mistakes ------------------------------------------------
    s = base_slide(prs, 6)
    add_heading(s, "Common mistakes", lead="A few pitfalls to avoid when first using Codex.")
    add_bullets(
        s,
        [
            "Overloading the prompt instead of moving durable rules into AGENTS.md or a skill",
            "Not telling the agent how to run build and test commands",
            "Skipping planning on multi-step, complex tasks",
            "Granting full machine permission before you understand the workflow",
            "Running live threads on the same files without Git worktrees",
            "Automating a recurring task before it is reliable manually",
            "Babysitting Codex instead of running it in parallel with your own work",
            "One thread per project (bloated context) instead of one thread per task",
        ],
        top=2.75,
        size=15.5,
        gap=7,
    )

    # --- 7. Deep research with subagents -----------------------------------
    s = base_slide(prs, 7)
    add_heading(s, "Deep research with subagents", lead="Codex has no separate “Deep Research” mode — build one by spawning parallel subagents.")
    add_bullets(
        s,
        [
            ("Live search — ", "run codex --search in the terminal (not inside a running chat)"),
            ("Ask explicitly — ", "“spawn N subagents in parallel” + “wait for all agents, then synthesize”"),
            ("Bounded roles — ", "primary-source · independent-source · skeptical reviewer · local-impact"),
            ("Each agent — ", "5–10 findings, source URLs + dates, facts vs. inferences, flag uncertainty"),
            ("Limits — ", "one coordinator + 3–5 agents; max_threads = 5, max_depth = 1"),
            ("Safety — ", "web content is untrusted; prefer read-only agents & primary sources"),
        ],
        top=2.75,
        size=15.5,
        gap=8,
    )

    # --- 8. Superpowers workflow -------------------------------------------
    s = base_slide(prs, 8)
    add_heading(s, "Superpowers solo workflow", lead="Ambiguous request → reviewed spec → TDD implementation → verified release.")
    add_bullets(
        s,
        [
            ("brainstorming — ", "approved spec with scope & acceptance criteria"),
            ("writing-plans — ", "task-sized TDD plan with exact files, commands, commits"),
            ("subagent-driven / executing-plans — ", "execute the plan task by task"),
            ("test-driven-development — ", "red → green → refactor for every behavior"),
            ("systematic-debugging — ", "reproduce, find root cause, add a regression test"),
            ("requesting / receiving-code-review — ", "severity-ranked review against the spec"),
            ("verification-before-completion — ", "fresh evidence, then commit & push main"),
        ],
        top=2.75,
        size=15.5,
        gap=7,
    )

    # --- 9. Thank you / resources ------------------------------------------
    s = base_slide(prs, 9)
    add_tag(s, "Thank you", MARGIN, 2.4)
    box = s.shapes.add_textbox(Inches(MARGIN), Inches(2.85), Inches(CONTENT_W), Inches(1.0))
    tf = box.text_frame
    _no_autofit(tf)
    tf.margin_left = tf.margin_top = 0
    p = tf.paragraphs[0]
    _set_run(p.add_run(), "Codex Community Build — Vienna", size=38, color=INK, font=FONT_SEMI)
    res = [
        ("github.com/JuliusGruber/codex-community-build-vienna", "https://github.com/JuliusGruber/codex-community-build-vienna"),
        ("github.com/JuliusGruber/AI-Engineering-Coach", "https://github.com/JuliusGruber/AI-Engineering-Coach"),
        ("codex-events.com/events/codex-build-vienna-2026-06-20", "https://codex-events.com/events/codex-build-vienna-2026-06-20?tab=details"),
    ]
    y = 4.25
    for label, url in res:
        lb = s.shapes.add_textbox(Inches(MARGIN), Inches(y), Inches(CONTENT_W), Inches(0.45))
        ltf = lb.text_frame
        _no_autofit(ltf)
        ltf.margin_left = ltf.margin_top = 0
        lp = ltf.paragraphs[0]
        _set_run(lp.add_run(), "→  ", size=16, color=INK, font=FONT_SEMI)
        run = lp.add_run()
        run.hyperlink.address = url  # set link first so the explicit color/underline below win
        _set_run(run, label, size=16, color=INK, font=FONT)
        run.font.underline = False
        y += 0.5

    set_link_theme_color(prs)

    out = Path(__file__).resolve().parent / "codex-community-build-vienna.pptx"
    prs.save(out)
    print(f"Saved {out} with {len(prs.slides._sldIdLst)} slides")


if __name__ == "__main__":
    build()
